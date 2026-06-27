import pandas as pd
from datetime import date, timedelta
import io
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

CORES_LOCAIS = ["AED6F1","A9DFBF","FAD7A0","D7BDE2","F9E79F","FADBD8"]
COR_HEADER = "1F4E79"
COR_HEADER2 = "2E75B6"
COR_FDS = "C5CAD6"
COR_FERIADO = "FFCCCC"
COR_BLOQUEIO = "FFE5CC"

def thin_border():
    s = Side(style="thin", color="BFBFBF")
    return Border(left=s, right=s, top=s, bottom=s)

def cs(ws, row, col, value, bold=False, bg=None, fc="000000", align="center", wrap=False, sz=9):
    c = ws.cell(row=row, column=col, value=value)
    c.font = Font(bold=bold, color=fc, size=sz)
    c.alignment = Alignment(horizontal=align, vertical="center", wrap_text=wrap)
    c.border = thin_border()
    if bg:
        c.fill = PatternFill("solid", fgColor=bg)
    return c

def gerar_datas(data_inicio_str, num_semanas):
    d0 = date.fromisoformat(data_inicio_str)
    return [[d0 + timedelta(weeks=s, days=d) for d in range(7)] for s in range(num_semanas)]

def montar_rodizio(num_semanas, num_sg, num_locais, ajustes_sg=None):
    rodizio = {}
    for sem in range(num_semanas):
        rodizio[sem] = {}
        for sg in range(1, num_sg + 1):
            local_idx = (sg - 1 + sem) % num_locais
            rodizio[sem][sg] = local_idx
    return rodizio

def gerar_escala(config):
    especialidade = config["especialidade"]
    grupo = config["grupo"]
    turma = config["turma"]
    data_inicio = config["data_inicio"]
    num_semanas = config["num_semanas"]
    num_sg = config["num_sg"]
    locais = config["locais"]
    alunos = config["alunos"]
    bloqueios = config["bloqueios"]
    pares_bp = config.get("pares_bp", {})
    ajustes = config.get("ajustes", [])

    feriados = set()
    for f in bloqueios.get("feriados", []):
        try:
            feriados.add(date.fromisoformat(f))
        except:
            pass

    semanas = gerar_datas(data_inicio, num_semanas)
    rodizio = montar_rodizio(num_semanas, num_sg, len(locais))

    # Aplicar ajustes de troca de SG
    alunos_por_sg_base = {}
    for a in alunos:
        sg = int(a["Sub Grupo"])
        alunos_por_sg_base.setdefault(sg, []).append(a)

    # alunos_por_sg_semana[sem][sg] = lista de alunos
    alunos_por_sg_semana = {}
    for sem in range(num_semanas):
        alunos_por_sg_semana[sem] = {}
        for sg in range(1, num_sg + 1):
            alunos_por_sg_semana[sem][sg] = list(alunos_por_sg_base.get(sg, []))

    for aj in ajustes:
        if aj["tipo"] == "Trocar de SG":
            sem = aj["semana"] - 1
            novo_sg = aj.get("novo_sg", 1)
            nome = aj["aluno"]
            # Encontrar aluno e mover
            for sg in range(1, num_sg + 1):
                for a in list(alunos_por_sg_semana[sem].get(sg, [])):
                    if a["Nome Completo"] == nome:
                        alunos_por_sg_semana[sem][sg].remove(a)
                        alunos_por_sg_semana[sem].setdefault(novo_sg, []).append(a)

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    _aba_subgrupos(wb, alunos_por_sg_base, num_sg)
    _aba_calendario(wb, semanas, rodizio, locais, num_sg, pares_bp, feriados)
    _aba_nominal(wb, semanas, rodizio, locais, alunos_por_sg_semana, num_sg, bloqueios, pares_bp, feriados)
    _aba_resumo(wb, semanas, alunos_por_sg_semana, num_sg, locais, rodizio, bloqueios, pares_bp, feriados)

    xlsx_io = io.BytesIO()
    wb.save(xlsx_io)
    xlsx_io.seek(0)

    pptx_io = _gerar_pptx(especialidade, grupo, turma, semanas, rodizio, locais, alunos_por_sg_base, num_sg, pares_bp)

    return {"xlsx": xlsx_io.read(), "pptx": pptx_io.read()}

def _aba_subgrupos(wb, alunos_por_sg, num_sg):
    ws = wb.create_sheet("Subgrupos")
    ws.sheet_view.showGridLines = False
    cs(ws,1,1,"SUBGRUPOS",bold=True,bg=COR_HEADER,fc="FFFFFF",sz=12)
    ws.merge_cells("A1:D1")
    ws.row_dimensions[1].height = 28
    for c,h in enumerate(["Sub Grupo","Nome Completo","RA","Turma"],1):
        cs(ws,2,c,h,bold=True,bg=COR_HEADER2,fc="FFFFFF")
    cores = ["D6E4F7","D6F7D6","F7F0D6","F7D6D6","F2D9F7","D9F7E8","FDE8D8","E8F8E8"]
    row = 3
    for sg in range(1, num_sg+1):
        cor = cores[(sg-1) % len(cores)]
        for a in alunos_por_sg.get(sg,[]):
            cs(ws,row,1,f"SG{sg}",bg=cor)
            cs(ws,row,2,a["Nome Completo"],bg=cor,align="left")
            cs(ws,row,3,str(a["RA"]),bg=cor)
            cs(ws,row,4,a.get("Turma","T6"),bg=cor)
            row += 1
    ws.column_dimensions["A"].width=10
    ws.column_dimensions["B"].width=40
    ws.column_dimensions["C"].width=15
    ws.column_dimensions["D"].width=8

def _aba_calendario(wb, semanas, rodizio, locais, num_sg, pares_bp, feriados):
    ws = wb.create_sheet("Calendário de Rodízio")
    ws.sheet_view.showGridLines = False
    cs(ws,1,1,"CALENDÁRIO DE RODÍZIO",bold=True,bg=COR_HEADER,fc="FFFFFF",sz=12)
    ws.merge_cells(f"A1:{get_column_letter(1+num_sg)}1")
    ws.row_dimensions[1].height=28
    headers = ["Semana / Período"] + [f"SG{i}" for i in range(1,num_sg+1)]
    for c,h in enumerate(headers,1):
        cs(ws,2,c,h,bold=True,bg=COR_HEADER2,fc="FFFFFF")
    cores_loc = {i: CORES_LOCAIS[i%len(CORES_LOCAIS)] for i in range(len(locais))}
    for sem_idx,semana in enumerate(semanas):
        row = sem_idx+3
        inicio = semana[0].strftime("%d/%m")
        fim = semana[-1].strftime("%d/%m")
        cs(ws,row,1,f"Sem {sem_idx+1} | {inicio}–{fim}",bold=True,bg="F2F2F2")
        par_bp = pares_bp.get(sem_idx,())
        for sg in range(1,num_sg+1):
            local_idx = rodizio[sem_idx].get(sg,0)
            local = locais[local_idx] if local_idx < len(locais) else {}
            nome = local.get("nome","?")
            cor = cores_loc.get(local_idx,"FFFFFF")
            if local.get("tipo") == "bloco_ped" and sg in par_bp:
                nome += " (ENF+PA)"
            cs(ws,row,sg+1,nome,bg=cor)
    ws.column_dimensions["A"].width=22
    for i in range(1,num_sg+1):
        ws.column_dimensions[get_column_letter(i+1)].width=22

def _aba_nominal(wb, semanas, rodizio, locais, alunos_por_sg_semana, num_sg, bloqueios, pares_bp, feriados):
    ws = wb.create_sheet("Escala Nominal Detalhada")
    ws.sheet_view.showGridLines = False
    dias_sem = ["Seg","Ter","Qua","Qui","Sex","Sáb","Dom"]
    total_cols = 2 + len(semanas)*7
    cs(ws,1,1,"ESCALA NOMINAL DETALHADA",bold=True,bg=COR_HEADER,fc="FFFFFF",sz=12)
    ws.merge_cells(f"A1:{get_column_letter(total_cols)}1")
    ws.row_dimensions[1].height=28
    cs(ws,2,1,"SG",bold=True,bg=COR_HEADER2,fc="FFFFFF")
    cs(ws,2,2,"Nome",bold=True,bg=COR_HEADER2,fc="FFFFFF")
    col = 3
    for sem_idx,semana in enumerate(semanas):
        inicio = semana[0].strftime("%d/%m")
        fim = semana[-1].strftime("%d/%m")
        cs(ws,2,col,f"Sem {sem_idx+1} | {inicio}–{fim}",bold=True,bg=COR_HEADER2,fc="FFFFFF")
        ws.merge_cells(start_row=2,start_column=col,end_row=2,end_column=col+6)
        for d,dia in enumerate(semana):
            cor_dia = COR_FERIADO if dia in feriados else "E8EDF5"
            cs(ws,3,col+d,f"{dias_sem[d]}\n{dia.strftime('%d/%m')}",bold=True,bg=cor_dia,wrap=True)
        col += 7
    ws.row_dimensions[3].height=30
    ws.freeze_panes = "C4"
    cores_sg = ["D6E4F7","D6F7D6","F7F0D6","F7D6D6","F2D9F7","D9F7E8","FDE8D8","E8F8E8"]
    row = 4
    for sg in range(1,num_sg+1):
        cor_sg = cores_sg[(sg-1)%len(cores_sg)]
        alunos_sg = set()
        for sem in range(len(semanas)):
            for a in alunos_por_sg_semana[sem].get(sg,[]):
                alunos_sg.add(a["Nome Completo"])
        for nome_aluno in sorted(alunos_sg):
            cs(ws,row,1,f"SG{sg}",bg=cor_sg,bold=True)
            cs(ws,row,2,nome_aluno,bg=cor_sg,align="left")
            col = 3
            for sem_idx,semana in enumerate(semanas):
                local_idx = rodizio[sem_idx].get(sg,0)
                local = locais[local_idx] if local_idx < len(locais) else {}
                nome_local = local.get("nome","?")
                cor_local = CORES_LOCAIS[local_idx%len(CORES_LOCAIS)]
                par_bp = pares_bp.get(sem_idx,())
                eh_bp = local.get("tipo") == "bloco_ped"
                for d,dia in enumerate(semana):
                    dia_sem = dia.weekday()
                    eh_fds = dia_sem >= 5
                    if dia in feriados:
                        cs(ws,row,col+d,"FERIADO",bg=COR_FERIADO)
                    elif eh_fds and not local.get("fds"):
                        cs(ws,row,col+d,"—",bg=COR_FDS)
                    elif eh_bp and sg in par_bp:
                        # No bloco pediátrico — mostra ENF+PA
                        if eh_fds:
                            cs(ws,row,col+d,"PA/M+T",bg=cor_local)
                        elif bloqueios.get("quinta_tarde") and dia_sem == 3:
                            cs(ws,row,col+d,"ENF/M\nPA/M",bg=cor_local,wrap=True)
                        else:
                            cs(ws,row,col+d,"ENF/M\nPA/M+T+C",bg=cor_local,wrap=True)
                    elif eh_bp:
                        # BP mas não é o par da vez — não vai ao BP
                        cs(ws,row,col+d,"—",bg="EEEEEE")
                    else:
                        turnos = []
                        if local.get("turno_m"): turnos.append("M")
                        if bloqueios.get("quinta_tarde") and dia_sem==3:
                            pass
                        elif bloqueios.get("terca_parcial") and dia_sem==1:
                            if local.get("turno_t"): turnos.append("T*")
                        else:
                            if local.get("turno_t"): turnos.append("T")
                        if local.get("turno_c") and (eh_fds or dia_sem in [4,5,6]):
                            turnos.append("C")
                        cs(ws,row,col+d,f"{nome_local}\n{'+'.join(turnos)}",bg=cor_local,wrap=True)
                col += 7
            row += 1
    ws.column_dimensions["A"].width=6
    ws.column_dimensions["B"].width=36
    for c in range(3,total_cols+1):
        ws.column_dimensions[get_column_letter(c)].width=11

def _aba_resumo(wb, semanas, alunos_por_sg_semana, num_sg, locais, rodizio, bloqueios, pares_bp, feriados):
    ws = wb.create_sheet("Resumo de Horas")
    ws.sheet_view.showGridLines = False
    n_sem = len(semanas)
    cs(ws,1,1,"RESUMO DE HORAS POR ALUNO",bold=True,bg=COR_HEADER,fc="FFFFFF",sz=12)
    ws.merge_cells(f"A1:{get_column_letter(5+n_sem)}1")
    ws.row_dimensions[1].height=28
    headers = ["SG","Nome","RA","Total Horas","Cinderelas"] + [f"Sem {s+1}" for s in range(n_sem)]
    for c,h in enumerate(headers,1):
        cs(ws,2,c,h,bold=True,bg=COR_HEADER2,fc="FFFFFF")
    cores_sg = ["D6E4F7","D6F7D6","F7F0D6","F7D6D6","F2D9F7","D9F7E8","FDE8D8","E8F8E8"]
    row = 3
    todos_alunos = {}
    for sg in range(1, num_sg+1):
        for sem in range(len(semanas)):
            for a in alunos_por_sg_semana[sem].get(sg,[]):
                if a["Nome Completo"] not in todos_alunos:
                    todos_alunos[a["Nome Completo"]] = {"sg": sg, "ra": a["RA"], "aluno": a}
    for sg in range(1, num_sg+1):
        cor_sg = cores_sg[(sg-1)%len(cores_sg)]
        alunos_sg = {n: d for n,d in todos_alunos.items() if d["sg"]==sg}
        for nome_aluno, info in sorted(alunos_sg.items()):
            horas_semanas = []
            total_cind = 0
            for sem_idx,semana in enumerate(semanas):
                local_idx = rodizio[sem_idx].get(sg,0)
                local = locais[local_idx] if local_idx < len(locais) else {}
                par_bp = pares_bp.get(sem_idx,())
                eh_bp = local.get("tipo") == "bloco_ped"
                horas = 0
                cind = 0
                for d,dia in enumerate(semana):
                    dia_sem = dia.weekday()
                    eh_fds = dia_sem >= 5
                    if dia in feriados:
                        continue
                    if eh_bp and sg in par_bp:
                        if eh_fds:
                            horas += 12  # PA sab/dom M+T
                        else:
                            horas += 6  # ENF manhã
                            if not (bloqueios.get("quinta_tarde") and dia_sem==3):
                                horas += 4  # PA cinderela dias úteis selecionados
                                if dia_sem in [0,2,4]: cind += 1
                    elif eh_bp:
                        pass
                    else:
                        if eh_fds and not local.get("fds"):
                            continue
                        if local.get("turno_m"): horas += 6
                        if bloqueios.get("quinta_tarde") and dia_sem==3:
                            pass
                        elif local.get("turno_t"):
                            if bloqueios.get("terca_parcial") and dia_sem==1:
                                horas += 3
                            else:
                                horas += 6
                        if local.get("turno_c") and (eh_fds or dia_sem in [4,5,6]):
                            horas += 4
                            cind += 1
                horas_semanas.append(horas)
                total_cind += cind
            cs(ws,row,1,f"SG{sg}",bg=cor_sg,bold=True)
            cs(ws,row,2,nome_aluno,bg=cor_sg,align="left")
            cs(ws,row,3,str(info["ra"]),bg=cor_sg)
            cs(ws,row,4,sum(horas_semanas),bg=cor_sg,bold=True)
            cs(ws,row,5,total_cind,bg=cor_sg)
            for s,h in enumerate(horas_semanas):
                cs(ws,row,6+s,h,bg=cor_sg)
            row += 1
    ws.column_dimensions["A"].width=6
    ws.column_dimensions["B"].width=36
    ws.column_dimensions["C"].width=14
    ws.column_dimensions["D"].width=14
    ws.column_dimensions["E"].width=12
    for s in range(n_sem):
        ws.column_dimensions[get_column_letter(6+s)].width=10

def _gerar_pptx(especialidade, grupo, turma, semanas, rodizio, locais, alunos_por_sg, num_sg, pares_bp):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def add_tb(slide, text, left, top, width, height, bold=False, size=18,
                   color=(0,0,0), bg=None, align=PP_ALIGN.CENTER):
            txBox = slide.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.bold = bold
            run.font.size = Pt(size)
            run.font.color.rgb = RGBColor(*color)
            if bg:
                txBox.fill.solid()
                txBox.fill.fore_color.rgb = RGBColor(*bg)
            return txBox

        # Capa
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(31,78,121)
        add_tb(slide,f"Escala de Internato\n{especialidade}",1,2,11,1.5,bold=True,size=32,color=(212,175,55))
        add_tb(slide,f"{grupo} | {turma}",1,4,11,1,size=20,color=(255,255,255))
        add_tb(slide,f"{semanas[0][0].strftime('%d/%m/%Y')} a {semanas[-1][-1].strftime('%d/%m/%Y')}",1,5,11,.8,size=16,color=(200,200,200))

        # Subgrupos
        slide = prs.slides.add_slide(blank)
        add_tb(slide,"SUBGRUPOS",0.3,0.2,12,.6,bold=True,size=22,color=(255,255,255),bg=(31,78,121))
        cores_sg = [(214,228,247),(214,247,214),(247,240,214),(247,214,214),(242,217,247),(217,247,232)]
        col_w = 13.33/min(num_sg,4)
        for sg in range(1,num_sg+1):
            col_idx=(sg-1)%4; row_idx=(sg-1)//4
            cor=cores_sg[(sg-1)%len(cores_sg)]
            left=0.1+col_idx*col_w; top=1.0+row_idx*3.0
            nomes="\n".join([a["Nome Completo"] for a in alunos_por_sg.get(sg,[])])
            add_tb(slide,f"SG{sg}\n{nomes}",left,top,col_w-.2,2.8,size=10,bg=cor,align=PP_ALIGN.LEFT)

        # Slides por semana
        cores_loc = [(174,214,241),(169,223,191),(250,215,160),(215,189,226),(249,231,159),(250,219,216)]
        for sem_idx,semana in enumerate(semanas):
            slide = prs.slides.add_slide(blank)
            inicio=semana[0].strftime("%d/%m"); fim=semana[-1].strftime("%d/%m")
            par_bp = pares_bp.get(sem_idx,())
            add_tb(slide,f"Semana {sem_idx+1} | {inicio}–{fim}",0.3,0.1,12,.55,bold=True,size=18,color=(255,255,255),bg=(31,78,121))
            n_locais=len(locais); col_w2=13.0/n_locais
            for loc_idx,local in enumerate(locais):
                cor=cores_loc[loc_idx%len(cores_loc)]
                left=0.15+loc_idx*col_w2
                titulo_local = local["nome"]
                if local.get("tipo")=="bloco_ped":
                    titulo_local += " (ENF+PA)"
                add_tb(slide,f"📍 {titulo_local}",left,.75,col_w2-.1,.45,bold=True,size=12,bg=cor)
                sgs_aqui=[sg for sg in range(1,num_sg+1) if rodizio[sem_idx].get(sg)==loc_idx]
                texto=""
                for sg in sgs_aqui:
                    em_bp = local.get("tipo")=="bloco_ped" and sg in par_bp
                    label = f"SG{sg}" + (" ★" if em_bp else "")
                    texto += f"{label}:\n"
                    for a in alunos_por_sg.get(sg,[]):
                        texto += f"  {a['Nome Completo']}\n"
                add_tb(slide,texto.strip(),left,1.25,col_w2-.1,5.8,size=8,bg=tuple(max(0,c-20) for c in cor),align=PP_ALIGN.LEFT)

        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io
    except ImportError:
        return io.BytesIO(b"")
